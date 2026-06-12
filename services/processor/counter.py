"""Accurate document page/character counting for all supported formats.

Uses pdfplumber for PDF (better accuracy than PyPDF2),
python-docx for DOCX, openpyxl for XLSX, python-pptx for PPTX,
and document_engine for everything else.
"""

import io
import math
import os
import chardet
import logging

logger = logging.getLogger(__name__)

CHARS_PER_PAGE = 1800
PRICE_PER_PAGE_TEXT = 68  # cents (€0.68 per 1800 chars — simple text/docx without glossary)
PRICE_PER_PAGE_DOC = 135  # cents (€1.35 per 1800 chars — complex/OCR formats)


def visible_char_count(text: str, is_html: bool = False) -> int:
    """Count billable characters close to what document translation APIs see.

    We keep spaces, tabs, and newlines instead of collapsing whitespace. This
    avoids undercounting formatted documents where paragraph/table separators
    are still processed by the translation provider. Only non-printing control
    characters are removed.
    """
    if not text:
        return 0

    if is_html:
        try:
            from bs4 import BeautifulSoup
            text = BeautifulSoup(text, "lxml").get_text(separator="\n", strip=True)
        except Exception:
            pass

    if text.startswith("\ufeff"):
        text = text[1:]

    text = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    count = 0
    for ch in text:
        code = ord(ch)
        if code < 32 and ch not in ("\n", "\t", "\f", "\v"):
            continue
        count += 1
    return count


def count_file(file_bytes: bytes, filename: str) -> dict:
    """Count pages and characters for a file. Returns pricing info."""
    ext = os.path.splitext(filename.lower())[1]

    try:
        if ext == ".pdf":
            return _count_pdf(file_bytes, filename)
        elif ext == ".docx":
            return _count_docx(file_bytes, filename)
        elif ext == ".doc":
            return _count_doc(file_bytes, filename)
        elif ext == ".xlsx":
            return _count_xlsx(file_bytes, filename)
        elif ext == ".xls":
            return _count_xls(file_bytes, filename)
        elif ext == ".pptx":
            return _count_pptx(file_bytes, filename)
        elif ext == ".ppt":
            return _count_ppt(file_bytes, filename)
        elif ext in (".jpg", ".jpeg", ".png", ".webp", ".svg"):
            return _count_image(file_bytes, filename, ext)
        elif ext == ".rtf":
            return _count_rtf(file_bytes, filename)
        elif ext in (".odt", ".ods", ".odp"):
            return _count_odf(file_bytes, filename, ext)
        elif ext in (".html", ".htm"):
            return _count_html(file_bytes, filename)
        elif ext == ".epub":
            return _count_epub(file_bytes, filename)
        elif ext in (
            ".txt", ".csv", ".tsv", ".md", ".xml", ".json", ".resjson",
            ".srt", ".sub", ".ass", ".ssa", ".vtt", ".po", ".xlf", ".xliff",
            ".go", ".yml", ".yaml", ".php", ".plist", ".stringsdict", ".tex",
        ):
            return _count_text(file_bytes, filename)
        elif ext in (".wps", ".et", ".dps", ".ai", ".indd", ".idml", ".chm", ".arxiv"):
            return _count_complex_document(file_bytes, filename, ext)
        else:
            return _count_text(file_bytes, filename)
    except Exception as e:
        logger.error(f"count_file failed for {filename}: {e}", exc_info=True)
        return {
            "pages": max(1, len(file_bytes) // 3000),
            "chars": 0,
            "file_type": ext.lstrip(".") or "unknown",
            "pricing_cents": max(1, len(file_bytes) // 3000) * PRICE_PER_PAGE_DOC,
            "method": "fallback",
            "error": str(e),
        }


def _count_pdf(file_bytes: bytes, filename: str) -> dict:
    """PDF: pdfplumber (primary) → PyPDF2 (fallback)."""
    total_chars = 0
    page_texts = []
    physical_pages = 0
    method = "pdf_exact"

    # Method 1: pdfplumber (best quality)
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            physical_pages = len(pdf.pages)
            for page in pdf.pages:
                text = page.extract_text() or ""
                page_texts.append(text)
            total_chars = visible_char_count("\n".join(page_texts))
            method = "pdfplumber"
    except Exception as e:
        logger.warning(f"pdfplumber count failed: {e}")

        # Method 2: PyPDF2 fallback
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes))
            physical_pages = len(reader.pages)
            page_texts = []
            for page in reader.pages:
                text = page.extract_text() or ""
                page_texts.append(text)
            total_chars = visible_char_count("\n".join(page_texts))
            method = "pypdf2"
        except Exception as e2:
            logger.warning(f"PyPDF2 count also failed: {e2}")
            # Last resort: physical pages only
            physical_pages = max(1, physical_pages)
            method = "pdf_page_count"

    # If text extraction yielded very little, use the full document engine with OCR fallback.
    if total_chars < physical_pages * 100:
        try:
            import document_engine
            result = document_engine.extract_text(file_bytes, filename)
            extracted = result.get("text", "")
            extracted_chars = visible_char_count(
                extracted,
                result.get("content_type") == "html",
            )
            if extracted_chars > total_chars:
                total_chars = extracted_chars
                method = result.get("method", "document_engine")
        except Exception as e:
            logger.warning(f"document_engine OCR count failed: {e}")

    if total_chars < physical_pages * 100:
        return {
            "pages": max(1, physical_pages),
            "chars": total_chars,
            "file_type": "pdf",
            "pricing_cents": max(1, physical_pages) * PRICE_PER_PAGE_DOC,
            "method": method,
            "note": "scanned_or_image_pdf",
        }

    char_pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
    return {
        "pages": char_pages,
        "chars": total_chars,
        "file_type": "pdf",
        "pricing_cents": char_pages * PRICE_PER_PAGE_DOC,
        "method": method,
    }


def _count_docx(file_bytes: bytes, filename: str) -> dict:
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))

    parts = []
    for para in doc.paragraphs:
        parts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)

    for section in doc.sections:
        for header_footer in [section.header, section.footer]:
            if header_footer is not None:
                for para in header_footer.paragraphs:
                    parts.append(para.text)

    total_chars = visible_char_count("\n".join(parts))
    pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
    return {
        "pages": pages,
        "chars": total_chars,
        "file_type": "docx",
        "pricing_cents": pages * PRICE_PER_PAGE_TEXT,
        "method": "docx_exact",
    }


def _count_doc(file_bytes: bytes, filename: str) -> dict:
    """Legacy .doc: try document_engine extraction, fallback to estimate."""
    try:
        import document_engine
        result = document_engine.extract_text(file_bytes, filename)
        text = result.get("text", "")
        if text.strip():
            total_chars = visible_char_count(text)
            pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
            return {
                "pages": pages,
                "chars": total_chars,
                "file_type": "doc",
                "pricing_cents": pages * PRICE_PER_PAGE_DOC,
                "method": "doc_extracted",
            }
    except Exception as e:
        logger.warning(f"doc extraction for counting failed: {e}")

    estimated_chars = len(file_bytes) // 3
    pages = max(1, math.ceil(estimated_chars / CHARS_PER_PAGE))
    return {
        "pages": pages,
        "chars": estimated_chars,
        "file_type": "doc",
        "pricing_cents": pages * PRICE_PER_PAGE_DOC,
        "method": "doc_estimate",
        "note": "legacy_format_estimated",
    }


def _count_xlsx(file_bytes: bytes, filename: str) -> dict:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)

    parts = []
    sheet_count = len(wb.sheetnames)

    for ws in wb.worksheets:
        parts.append(ws.title)
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    parts.append(str(cell.value))

    wb.close()

    total_chars = visible_char_count("\n".join(parts))
    pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
    return {
        "pages": pages,
        "chars": total_chars,
        "file_type": "xlsx",
        "sheets": sheet_count,
        "pricing_cents": pages * PRICE_PER_PAGE_TEXT,
        "method": "xlsx_exact",
    }


def _count_xls(file_bytes: bytes, filename: str) -> dict:
    """Legacy .xls: try LibreOffice conversion → count as xlsx."""
    try:
        import document_engine
        result = document_engine.extract_text(file_bytes, filename)
        text = result.get("text", "")
        if text.strip():
            total_chars = visible_char_count(text)
            pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
            return {
                "pages": pages,
                "chars": total_chars,
                "file_type": "xls",
                "pricing_cents": pages * PRICE_PER_PAGE_TEXT,
                "method": "xls_extracted",
            }
    except Exception:
        pass

    estimated_chars = len(file_bytes) // 4
    pages = max(1, math.ceil(estimated_chars / CHARS_PER_PAGE))
    return {
        "pages": pages,
        "chars": estimated_chars,
        "file_type": "xls",
        "pricing_cents": pages * PRICE_PER_PAGE_DOC,
        "method": "xls_estimate",
    }


def _count_pptx(file_bytes: bytes, filename: str) -> dict:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slide_count = len(prs.slides)

    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            parts.append(slide.notes_slide.notes_text_frame.text)

    total_chars = visible_char_count("\n".join(parts))
    char_pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
    return {
        "pages": char_pages,
        "chars": total_chars,
        "slides": slide_count,
        "file_type": "pptx",
        "pricing_cents": char_pages * PRICE_PER_PAGE_DOC,
        "method": "pptx_exact",
    }


def _count_ppt(file_bytes: bytes, filename: str) -> dict:
    """Legacy PPT: extract text when possible, otherwise count as page/slide estimate."""
    try:
        import document_engine
        result = document_engine.extract_text(file_bytes, filename)
        text = result.get("text", "")
        if text.strip():
            total_chars = visible_char_count(text, result.get("content_type") == "html")
            pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
            return {
                "pages": pages,
                "chars": total_chars,
                "file_type": "ppt",
                "pricing_cents": pages * PRICE_PER_PAGE_DOC,
                "method": "ppt_extracted",
            }
    except Exception as e:
        logger.warning(f"ppt extraction for counting failed: {e}")

    pages = max(1, math.ceil(len(file_bytes) / 15000))
    return {
        "pages": pages,
        "chars": 0,
        "file_type": "ppt",
        "pricing_cents": pages * PRICE_PER_PAGE_DOC,
        "method": "ppt_estimate",
    }


def _count_image(file_bytes: bytes, filename: str, ext: str) -> dict:
    """Images are translated/OCRed as visual documents; bill at least one page."""
    return {
        "pages": 1,
        "chars": 0,
        "file_type": ext.lstrip("."),
        "pricing_cents": PRICE_PER_PAGE_DOC,
        "method": "image_page_count",
        "note": "image_or_scan",
    }


def _count_rtf(file_bytes: bytes, filename: str) -> dict:
    """RTF: striprtf for text extraction."""
    try:
        from striprtf.striprtf import rtf_to_text
        text = rtf_to_text(file_bytes.decode("utf-8", errors="replace"))
        total_chars = visible_char_count(text)
        pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
        return {
            "pages": pages,
            "chars": total_chars,
            "file_type": "rtf",
            "pricing_cents": pages * PRICE_PER_PAGE_TEXT,
            "method": "rtf_exact",
        }
    except Exception as e:
        logger.warning(f"RTF count failed: {e}")
        estimated = len(file_bytes) // 3
        pages = max(1, math.ceil(estimated / CHARS_PER_PAGE))
        return {
            "pages": pages,
            "chars": estimated,
            "file_type": "rtf",
            "pricing_cents": pages * PRICE_PER_PAGE_TEXT,
            "method": "rtf_estimate",
        }


def _count_odf(file_bytes: bytes, filename: str, ext: str) -> dict:
    """ODF formats: use document_engine extraction."""
    try:
        import document_engine
        result = document_engine.extract_text(file_bytes, filename)
        text = result.get("text", "")
        if text.strip():
            total_chars = visible_char_count(text, result.get("content_type") == "html")
            pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
            price = PRICE_PER_PAGE_TEXT if ext == ".odt" else PRICE_PER_PAGE_DOC
            return {
                "pages": pages,
                "chars": total_chars,
                "file_type": ext.lstrip("."),
                "pricing_cents": pages * price,
                "method": f"{ext.lstrip('.')}_extracted",
            }
    except Exception:
        pass

    estimated = len(file_bytes) // 3
    pages = max(1, math.ceil(estimated / CHARS_PER_PAGE))
    return {
        "pages": pages,
        "chars": estimated,
        "file_type": ext.lstrip("."),
        "pricing_cents": pages * PRICE_PER_PAGE_DOC,
        "method": f"{ext.lstrip('.')}_estimate",
    }


def _count_html(file_bytes: bytes, filename: str) -> dict:
    """HTML: BeautifulSoup text extraction for accurate counting."""
    try:
        from bs4 import BeautifulSoup
        detected = chardet.detect(file_bytes)
        encoding = detected.get("encoding") or "utf-8"
        try:
            html = file_bytes.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            html = file_bytes.decode("utf-8", errors="replace")

        soup = BeautifulSoup(html, "lxml")
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        total_chars = visible_char_count(text)
        pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
        return {
            "pages": pages,
            "chars": total_chars,
            "file_type": "html",
            "pricing_cents": pages * PRICE_PER_PAGE_TEXT,
            "method": "html_exact",
        }
    except Exception:
        return _count_text(file_bytes, filename)


def _count_epub(file_bytes: bytes, filename: str) -> dict:
    """EPUB: ebooklib extraction for counting."""
    try:
        import document_engine
        result = document_engine.extract_text(file_bytes, filename)
        text = result.get("text", "")
        if text.strip():
            total_chars = visible_char_count(text, result.get("content_type") == "html")
            pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
            return {
                "pages": pages,
                "chars": total_chars,
                "file_type": "epub",
                "pricing_cents": pages * PRICE_PER_PAGE_DOC,
                "method": "epub_exact",
            }
    except Exception:
        pass

    estimated = len(file_bytes) // 4
    pages = max(1, math.ceil(estimated / CHARS_PER_PAGE))
    return {
        "pages": pages,
        "chars": estimated,
        "file_type": "epub",
        "pricing_cents": pages * PRICE_PER_PAGE_DOC,
        "method": "epub_estimate",
    }


def _count_complex_document(file_bytes: bytes, filename: str, ext: str) -> dict:
    """Provider-supported complex formats: extract when possible, else conservative estimate."""
    try:
        import document_engine
        result = document_engine.extract_text(file_bytes, filename)
        text = result.get("text", "")
        if text.strip():
            total_chars = visible_char_count(text, result.get("content_type") == "html")
            pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))
            return {
                "pages": pages,
                "chars": total_chars,
                "file_type": ext.lstrip("."),
                "pricing_cents": pages * PRICE_PER_PAGE_DOC,
                "method": f"{ext.lstrip('.')}_extracted",
            }
    except Exception as e:
        logger.warning(f"{ext} extraction for counting failed: {e}")

    estimated_chars = max(CHARS_PER_PAGE, len(file_bytes) // 3)
    pages = max(1, math.ceil(estimated_chars / CHARS_PER_PAGE))
    return {
        "pages": pages,
        "chars": estimated_chars,
        "file_type": ext.lstrip("."),
        "pricing_cents": pages * PRICE_PER_PAGE_DOC,
        "method": f"{ext.lstrip('.')}_estimate",
        "note": "complex_format_estimated",
    }


def _count_text(file_bytes: bytes, filename: str) -> dict:
    ext = os.path.splitext(filename.lower())[1].lstrip(".")

    detected = chardet.detect(file_bytes)
    encoding = detected.get("encoding") or "utf-8"

    try:
        text = file_bytes.decode(encoding)
    except (UnicodeDecodeError, LookupError):
        text = file_bytes.decode("utf-8", errors="replace")

    if text.startswith("\ufeff"):
        text = text[1:]

    total_chars = visible_char_count(text)
    pages = max(1, math.ceil(total_chars / CHARS_PER_PAGE))

    return {
        "pages": pages,
        "chars": total_chars,
        "file_type": ext or "txt",
        "pricing_cents": pages * PRICE_PER_PAGE_TEXT,
        "method": "text_exact",
    }
