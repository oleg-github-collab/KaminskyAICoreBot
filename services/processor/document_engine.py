"""Document Engine — unified text extraction for all formats.

Extraction priority per format:
  PDF:  pdfplumber → PyPDF2 → OCR (tesseract) → pdftotext
  DOCX: python-docx (paragraphs + tables + headers)
  DOC:  LibreOffice → mammoth → size heuristic
  XLSX: openpyxl (all sheets, all cells)
  PPTX: python-pptx (slides + notes + tables)
  RTF:  striprtf
  ODT:  python-docx via odt→docx or direct XML
  HTML: BeautifulSoup
  EPUB: ebooklib + BeautifulSoup
  TXT/CSV/JSON/XML: chardet decode
  Images: OCR (tesseract)
"""

import io
import os
import logging
import subprocess
import tempfile
from typing import Optional

logger = logging.getLogger(__name__)

MAX_TEXT_LENGTH = 500_000  # 500k chars max


def extract_text(file_bytes: bytes, filename: str) -> dict:
    """Extract text from any supported document format.

    Returns: {text, length, method, format, pages?, metadata?}
    """
    ext = os.path.splitext(filename.lower())[1].lstrip(".")

    try:
        if ext == "pdf":
            return _extract_pdf(file_bytes, filename)
        elif ext == "docx":
            return _extract_docx(file_bytes, filename)
        elif ext == "doc":
            return _extract_doc(file_bytes, filename)
        elif ext == "xlsx":
            return _extract_xlsx(file_bytes, filename)
        elif ext == "xls":
            return _extract_xls(file_bytes, filename)
        elif ext == "pptx":
            return _extract_pptx(file_bytes, filename)
        elif ext == "rtf":
            return _extract_rtf(file_bytes, filename)
        elif ext in ("odt", "ods", "odp"):
            return _extract_odf(file_bytes, filename, ext)
        elif ext in ("html", "htm"):
            return _extract_html(file_bytes, filename)
        elif ext == "epub":
            return _extract_epub(file_bytes, filename)
        elif ext in ("png", "jpg", "jpeg", "tiff", "tif", "bmp", "webp"):
            return _extract_image_ocr(file_bytes, filename)
        else:
            return _extract_text_file(file_bytes, filename)
    except Exception as e:
        logger.error(f"Extraction failed for {filename}: {e}", exc_info=True)
        # Last resort: try as plain text
        try:
            text = file_bytes.decode("utf-8", errors="replace")
            return _result(text, "fallback_decode", ext or "unknown")
        except Exception:
            return {"text": "", "length": 0, "method": "failed", "format": ext, "error": str(e)}


def extract_metadata(file_bytes: bytes, filename: str) -> dict:
    """Extract document metadata (author, title, dates, pages, language)."""
    ext = os.path.splitext(filename.lower())[1].lstrip(".")
    meta = {"filename": filename, "format": ext, "size_bytes": len(file_bytes)}

    try:
        if ext == "pdf":
            import pdfplumber
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                meta["pages"] = len(pdf.pages)
                info = pdf.metadata or {}
                meta["title"] = info.get("Title", "")
                meta["author"] = info.get("Author", "")
                meta["creator"] = info.get("Creator", "")
                meta["created"] = info.get("CreationDate", "")
                meta["modified"] = info.get("ModDate", "")

        elif ext == "docx":
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            props = doc.core_properties
            meta["title"] = props.title or ""
            meta["author"] = props.author or ""
            meta["created"] = str(props.created or "")
            meta["modified"] = str(props.modified or "")
            meta["paragraphs"] = len(doc.paragraphs)
            meta["tables"] = len(doc.tables)

        elif ext == "xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(file_bytes), read_only=True)
            meta["sheets"] = wb.sheetnames
            meta["sheet_count"] = len(wb.sheetnames)
            wb.close()

        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            meta["slides"] = len(prs.slides)

        # Detect language
        try:
            text_result = extract_text(file_bytes, filename)
            sample = text_result.get("text", "")[:2000]
            if sample.strip():
                from langdetect import detect
                meta["language"] = detect(sample)
        except Exception:
            pass

    except Exception as e:
        meta["error"] = str(e)

    return meta


def extract_tables(file_bytes: bytes, filename: str) -> dict:
    """Extract tables from PDF, DOCX, XLSX, PPTX."""
    ext = os.path.splitext(filename.lower())[1].lstrip(".")
    tables = []

    try:
        if ext == "pdf":
            import pdfplumber
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for i, page in enumerate(pdf.pages):
                    for table in page.extract_tables() or []:
                        if table:
                            tables.append({
                                "page": i + 1,
                                "rows": len(table),
                                "cols": max(len(r) for r in table) if table else 0,
                                "data": table[:50],  # Limit rows
                            })

        elif ext == "docx":
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            for i, table in enumerate(doc.tables):
                rows = []
                for row in table.rows[:50]:
                    rows.append([cell.text for cell in row.cells])
                tables.append({
                    "index": i,
                    "rows": len(table.rows),
                    "cols": len(table.columns),
                    "data": rows,
                })

        elif ext in ("xlsx", "xls"):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(max_row=100, values_only=True):
                    rows.append([str(c) if c is not None else "" for c in row])
                if rows:
                    tables.append({
                        "sheet": sheet_name,
                        "rows": len(rows),
                        "cols": max(len(r) for r in rows) if rows else 0,
                        "data": rows,
                    })
            wb.close()

        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.has_table:
                        rows = []
                        for row in shape.table.rows[:50]:
                            rows.append([cell.text for cell in row.cells])
                        tables.append({
                            "slide": slide_idx + 1,
                            "rows": len(shape.table.rows),
                            "cols": len(shape.table.columns),
                            "data": rows,
                        })

    except Exception as e:
        logger.error(f"Table extraction failed for {filename}: {e}")
        return {"tables": [], "error": str(e)}

    return {"tables": tables, "count": len(tables)}


# ─── Format-specific extractors ─────────────────────────────────────────────

def _extract_pdf(file_bytes: bytes, filename: str) -> dict:
    """PDF: PyMuPDF rich HTML → pdfplumber → PyPDF2 → pdftotext → OCR."""
    pages = 0

    # Primary: rich HTML extraction with PyMuPDF (font-based headings, bold/italic)
    try:
        html, pages = _extract_pdf_rich(file_bytes)
        if html and len(html.strip()) > 50:
            result = _result(html, "pdf_rich", "pdf", content_type="html")
            result["pages"] = pages
            return result
    except Exception as e:
        logger.warning(f"PyMuPDF rich extraction failed: {e}")

    # Fallback: plain text pipeline
    text = ""
    method = ""

    # Method 1: pdfplumber (best for layout-preserving extraction)
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            pages = len(pdf.pages)
            parts = []
            for page in pdf.pages:
                page_text = page.extract_text(layout=True) or ""
                parts.append(page_text)
            text = "\n\n".join(parts)
            method = "pdfplumber"
    except Exception as e:
        logger.warning(f"pdfplumber failed: {e}")

    # Method 2: PyPDF2 fallback
    if len(text.strip()) < 50 and not method.endswith("_ocr"):
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes))
            pages = len(reader.pages)
            parts = []
            for page in reader.pages:
                parts.append(page.extract_text() or "")
            pypdf_text = "\n\n".join(parts)
            if len(pypdf_text.strip()) > len(text.strip()):
                text = pypdf_text
                method = "pypdf2"
        except Exception as e:
            logger.warning(f"PyPDF2 failed: {e}")

    # Method 3: pdftotext (poppler) via subprocess
    if len(text.strip()) < 50:
        try:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            result = subprocess.run(
                ["pdftotext", "-layout", "-enc", "UTF-8", tmp_path, "-"],
                capture_output=True, text=True, timeout=60
            )
            os.unlink(tmp_path)
            if result.returncode == 0 and len(result.stdout.strip()) > len(text.strip()):
                text = result.stdout
                method = "pdftotext"
        except Exception as e:
            logger.warning(f"pdftotext failed: {e}")

    # Method 4: OCR for scanned PDFs
    if len(text.strip()) < pages * 50 and pages > 0:
        ocr_text = _ocr_pdf(file_bytes)
        if ocr_text and len(ocr_text.strip()) > len(text.strip()):
            text = ocr_text
            method = "ocr_tesseract"

    result = _result(text, method or "pdf_empty", "pdf")
    result["pages"] = pages
    return result


def _ocr_pdf(file_bytes: bytes) -> Optional[str]:
    """OCR a PDF using pdf2image + tesseract."""
    try:
        from pdf2image import convert_from_bytes
        import pytesseract

        images = convert_from_bytes(file_bytes, dpi=300, fmt="png")
        parts = []
        for img in images[:50]:  # Limit to 50 pages
            text = pytesseract.image_to_string(img, lang="ukr+deu+eng+rus+pol")
            parts.append(text)
        return "\n\n".join(parts)
    except Exception as e:
        logger.warning(f"OCR failed: {e}")
        return None


def _clean_pdf_text(text: str) -> str:
    """Clean PyMuPDF extracted text — fix encoding artifacts."""
    import re
    # Replace Unicode replacement character with space
    text = text.replace('\ufffd', ' ')
    # Replace various Unicode spaces with regular space
    text = re.sub(r'[\u00a0\u2000-\u200b\u202f\u205f\u3000\ufeff]', ' ', text)
    # Collapse multiple spaces
    text = re.sub(r'  +', ' ', text)
    return text


def _extract_pdf_rich(file_bytes: bytes) -> tuple:
    """Build semantic HTML from PDF using PyMuPDF font analysis.

    Detects headings by font size ratio, bold/italic by font flags,
    superscript, monospace. Returns (html_string, page_count) or (None, count).
    """
    import fitz  # PyMuPDF
    from html import escape as html_esc
    from collections import Counter

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    num_pages = len(doc)

    if num_pages == 0:
        doc.close()
        return None, 0

    # Pass 1: collect font size statistics (weighted by text length)
    size_chars = Counter()
    for page in doc:
        try:
            d = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
            for block in d.get("blocks", []):
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        t = span.get("text", "").strip()
                        if t:
                            size_chars[round(span.get("size", 12), 1)] += len(t)
        except Exception:
            continue

    if not size_chars:
        doc.close()
        return None, num_pages

    body_size = size_chars.most_common(1)[0][0]

    # Pass 2: build HTML
    html_parts = []

    for page_idx in range(num_pages):
        page = doc[page_idx]
        if page_idx > 0:
            html_parts.append('<hr class="pdf-page-break">')

        try:
            d = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
        except Exception:
            continue

        for block in d.get("blocks", []):
            if block.get("type") != 0:
                continue

            lines_data = []
            for line in block.get("lines", []):
                spans = []
                for span in line.get("spans", []):
                    text = _clean_pdf_text(span.get("text", ""))
                    if not text.strip():
                        if text:
                            spans.append({"t": " ", "s": span.get("size", body_size), "b": False, "i": False, "sup": False, "m": False})
                        continue
                    flags = span.get("flags", 0)
                    spans.append({
                        "t": text,
                        "s": span.get("size", body_size),
                        "b": bool(flags & 16),   # bold
                        "i": bool(flags & 2),    # italic
                        "sup": bool(flags & 1),  # superscript
                        "m": bool(flags & 8),    # monospace
                    })
                if spans:
                    lines_data.append(spans)

            if not lines_data:
                continue

            # Weighted average font size for this block
            total = sum(len(s["t"]) for ln in lines_data for s in ln)
            if total == 0:
                continue

            avg = sum(s["s"] * len(s["t"]) for ln in lines_data for s in ln) / total
            all_bold = all(s["b"] for ln in lines_data for s in ln if s["t"].strip())

            # Determine HTML tag by font size ratio
            r = avg / body_size if body_size > 0 else 1.0
            if r > 1.6:
                tag = "h1"
            elif r > 1.3:
                tag = "h2"
            elif r > 1.1 and all_bold:
                tag = "h3"
            elif all_bold and total < 200:
                tag = "h4"
            else:
                tag = "p"

            # Build inline HTML with formatting
            parts = []
            for li, line_spans in enumerate(lines_data):
                if li > 0:
                    parts.append("<br>" if tag == "p" else " ")
                for span in line_spans:
                    t = html_esc(span["t"])
                    if span["sup"]:
                        t = f"<sup>{t}</sup>"
                    if span["m"] and tag == "p":
                        t = f"<code>{t}</code>"
                    if span["b"] and span["i"]:
                        t = f"<strong><em>{t}</em></strong>" if tag == "p" else f"<em>{t}</em>"
                    elif span["b"] and tag == "p":
                        t = f"<strong>{t}</strong>"
                    elif span["i"]:
                        t = f"<em>{t}</em>"
                    parts.append(t)

            content = "".join(parts).strip()
            if content:
                html_parts.append(f"<{tag}>{content}</{tag}>")

    doc.close()

    html = "\n".join(html_parts)
    if not html.strip():
        return None, num_pages

    return html, num_pages


def _extract_docx(file_bytes: bytes, filename: str) -> dict:
    """DOCX: mammoth for rich semantic HTML, fallback to python-docx markdown."""
    # Primary: mammoth → clean semantic HTML (headings, bold, italic, tables, footnotes, lists)
    try:
        html = _extract_docx_mammoth(file_bytes)
        if html and len(html.strip()) > 20:
            return _result(html, "docx_mammoth", "docx", content_type="html")
    except Exception as e:
        logger.warning(f"mammoth DOCX extraction failed: {e}")

    # Fallback: python-docx with markdown formatting
    return _extract_docx_markdown(file_bytes, filename)


def _extract_docx_mammoth(file_bytes: bytes) -> str:
    """Convert DOCX to semantic HTML using mammoth.

    Handles: headings h1-h6, bold, italic, underline, strikethrough, lists (nested),
    tables, footnotes/endnotes, links, superscript/subscript.
    """
    import mammoth
    import re
    from html import escape as html_esc

    # Style map for custom/localized styles mammoth doesn't auto-detect
    style_map = "\n".join([
        "p[style-name='Title'] => h1:fresh",
        "p[style-name='Subtitle'] => h2:fresh",
        "p[style-name='Titel'] => h1:fresh",
        "p[style-name='Untertitel'] => h2:fresh",
        "p[style-name='Quote'] => blockquote:fresh",
        "p[style-name='Intense Quote'] => blockquote:fresh",
        "p[style-name='Zitat'] => blockquote:fresh",
        "p[style-name='Intensives Zitat'] => blockquote:fresh",
        "r[style-name='Intense Emphasis'] => strong > em:fresh",
    ])

    result = mammoth.convert_to_html(
        io.BytesIO(file_bytes),
        style_map=style_map,
    )

    html = result.value

    for msg in result.messages:
        logger.info(f"mammoth: {msg}")

    if not html or not html.strip():
        return None

    # Strip images (can be huge data URIs, not useful for text preview)
    html = re.sub(r'<img[^>]*/?\s*>', '', html)

    # Add headers/footers from python-docx (mammoth doesn't extract these)
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        hf_parts = []
        seen = set()
        for section in doc.sections:
            for hf in [section.header, section.footer]:
                if hf:
                    for para in hf.paragraphs:
                        text = para.text.strip()
                        if text and text not in seen:
                            hf_parts.append(text)
                            seen.add(text)
        if hf_parts:
            hf_html = '</p><p>'.join(html_esc(t) for t in hf_parts)
            html += f'\n<div class="doc-header-footer"><p>{hf_html}</p></div>'
    except Exception as e:
        logger.warning(f"Headers/footers extraction failed: {e}")

    return html


def _extract_docx_markdown(file_bytes: bytes, filename: str) -> dict:
    """DOCX fallback: python-docx with markdown formatting."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = (para.style.name or "").lower() if para.style else ""
        style_id = (para.style.style_id or "").lower() if para.style else ""
        heading_level = _detect_heading_level(style_name, style_id)

        formatted = _format_runs(para)
        if not formatted:
            continue

        if heading_level == 1:
            parts.append(f"# {formatted}")
        elif heading_level == 2:
            parts.append(f"## {formatted}")
        elif heading_level == 3:
            parts.append(f"### {formatted}")
        elif heading_level >= 4:
            parts.append(f"#### {formatted}")
        elif "list" in style_name or "bullet" in style_name or "liste" in style_name:
            parts.append(f"- {formatted}")
        elif "quote" in style_name or "zitat" in style_name:
            parts.append(f"> {formatted}")
        else:
            parts.append(formatted)

    for table in doc.tables:
        table_text = []
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                table_text.append(row_text)
        if table_text:
            parts.append("\n".join(table_text))

    for section in doc.sections:
        for hf in [section.header, section.footer]:
            if hf:
                for para in hf.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)

    return _result("\n\n".join(parts), "docx_markdown", "docx")


def _detect_heading_level(style_name: str, style_id: str) -> int:
    """Detect heading level from DOCX style name/id. Returns 0 for non-headings."""
    # Check style_id first (more reliable, standardized)
    if "heading1" in style_id or style_id == "title":
        return 1
    if "heading2" in style_id or style_id == "subtitle":
        return 2
    if "heading3" in style_id:
        return 3
    if "heading4" in style_id or "heading5" in style_id or "heading6" in style_id:
        return 4

    # Fallback: check localized style_name
    for pattern, level in [
        ("heading 1", 1), ("heading 2", 2), ("heading 3", 3), ("heading 4", 4),
        ("überschrift 1", 1), ("überschrift 2", 2), ("überschrift 3", 3), ("überschrift 4", 4),
        ("заголовок 1", 1), ("заголовок 2", 2), ("заголовок 3", 3), ("заголовок 4", 4),
        ("titre 1", 1), ("titre 2", 2), ("titre 3", 3),
        ("title", 1), ("subtitle", 2), ("titel", 1), ("untertitel", 2),
    ]:
        if pattern in style_name:
            return level

    # Generic heading detection
    if "heading" in style_id or "heading" in style_name:
        return 3

    return 0


def _format_runs(para) -> str:
    """Format paragraph runs with bold/italic as markdown markers."""
    if not para.runs:
        return para.text.strip()

    parts = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        is_bold = run.bold is True
        is_italic = run.italic is True

        if is_bold and is_italic:
            parts.append(f"***{text}***")
        elif is_bold:
            parts.append(f"**{text}**")
        elif is_italic:
            parts.append(f"*{text}*")
        else:
            parts.append(text)

    result = "".join(parts).strip()

    # Clean up adjacent identical markers: **text****more** → **text more**
    result = result.replace("****", " ").replace("****** ", " ")
    result = result.replace("** **", " ").replace("***  ***", " ")

    return result


def _extract_doc(file_bytes: bytes, filename: str) -> dict:
    """Legacy .doc: LibreOffice conversion → docx extraction."""
    # Try LibreOffice conversion
    try:
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        outdir = tempfile.mkdtemp()
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "docx", "--outdir", outdir, tmp_path],
            capture_output=True, text=True, timeout=30
        )
        os.unlink(tmp_path)

        if result.returncode == 0:
            docx_path = os.path.join(outdir, os.path.splitext(os.path.basename(tmp_path))[0] + ".docx")
            if os.path.exists(docx_path):
                with open(docx_path, "rb") as f:
                    docx_bytes = f.read()
                os.unlink(docx_path)
                os.rmdir(outdir)
                r = _extract_docx(docx_bytes, filename.replace(".doc", ".docx"))
                r["method"] = "doc_via_libreoffice"
                return r

        # Cleanup
        for f in os.listdir(outdir):
            os.unlink(os.path.join(outdir, f))
        os.rmdir(outdir)
    except Exception as e:
        logger.warning(f"LibreOffice conversion failed: {e}")

    # Fallback: mammoth (rich HTML conversion — works for .doc too)
    try:
        html = _extract_docx_mammoth(file_bytes)
        if html and len(html.strip()) > 20:
            return _result(html, "doc_mammoth", "doc", content_type="html")
    except Exception as e:
        logger.warning(f"mammoth failed: {e}")

    # Last resort: size heuristic
    estimated = len(file_bytes) // 3
    return {"text": "", "length": 0, "method": "doc_unsupported", "format": "doc",
            "note": f"Legacy .doc format, estimated ~{estimated} chars"}


def _extract_xlsx(file_bytes: bytes, filename: str) -> dict:
    """XLSX: all sheets, all cells, with sheet names."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_lines = [f"=== {sheet_name} ==="]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells).strip()
            if line and line != "|":
                sheet_lines.append(line)
        if len(sheet_lines) > 1:
            parts.append("\n".join(sheet_lines))

    wb.close()
    return _result("\n\n".join(parts), "xlsx_full", "xlsx")


def _extract_xls(file_bytes: bytes, filename: str) -> dict:
    """Legacy .xls: try LibreOffice conversion."""
    try:
        with tempfile.NamedTemporaryFile(suffix=".xls", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        outdir = tempfile.mkdtemp()
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "xlsx", "--outdir", outdir, tmp_path],
            capture_output=True, timeout=30
        )
        os.unlink(tmp_path)

        xlsx_path = os.path.join(outdir, os.path.splitext(os.path.basename(tmp_path))[0] + ".xlsx")
        if os.path.exists(xlsx_path):
            with open(xlsx_path, "rb") as f:
                xlsx_bytes = f.read()
            os.unlink(xlsx_path)
            os.rmdir(outdir)
            r = _extract_xlsx(xlsx_bytes, filename.replace(".xls", ".xlsx"))
            r["method"] = "xls_via_libreoffice"
            return r
    except Exception as e:
        logger.warning(f"xls conversion failed: {e}")

    return {"text": "", "length": 0, "method": "xls_unsupported", "format": "xls"}


def _extract_pptx(file_bytes: bytes, filename: str) -> dict:
    """PPTX: slides + speaker notes + tables."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    parts = []

    for i, slide in enumerate(prs.slides):
        slide_parts = [f"--- Slide {i + 1} ---"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        slide_parts.append(row_text)

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[Notes: {notes}]")

        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))

    result = _result("\n\n".join(parts), "pptx_full", "pptx")
    result["slides"] = len(prs.slides)
    return result


def _extract_rtf(file_bytes: bytes, filename: str) -> dict:
    """RTF: striprtf library."""
    try:
        from striprtf.striprtf import rtf_to_text
        text = rtf_to_text(file_bytes.decode("utf-8", errors="replace"))
        return _result(text, "rtf_striprtf", "rtf")
    except Exception as e:
        logger.warning(f"RTF extraction failed: {e}")
        return {"text": "", "length": 0, "method": "rtf_failed", "format": "rtf", "error": str(e)}


def _extract_odf(file_bytes: bytes, filename: str, ext: str) -> dict:
    """ODT/ODS/ODP: LibreOffice conversion → extract."""
    target_fmt = {"odt": "docx", "ods": "xlsx", "odp": "pptx"}.get(ext, "docx")
    extractor = {"docx": _extract_docx, "xlsx": _extract_xlsx, "pptx": _extract_pptx}.get(target_fmt)

    try:
        with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        outdir = tempfile.mkdtemp()
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", target_fmt, "--outdir", outdir, tmp_path],
            capture_output=True, timeout=30
        )
        os.unlink(tmp_path)

        out_path = os.path.join(outdir, os.path.splitext(os.path.basename(tmp_path))[0] + f".{target_fmt}")
        if os.path.exists(out_path) and extractor:
            with open(out_path, "rb") as f:
                converted_bytes = f.read()
            os.unlink(out_path)
            os.rmdir(outdir)
            r = extractor(converted_bytes, filename)
            r["method"] = f"{ext}_via_libreoffice"
            return r
    except Exception as e:
        logger.warning(f"ODF conversion failed: {e}")

    return {"text": "", "length": 0, "method": f"{ext}_failed", "format": ext}


def _extract_html(file_bytes: bytes, filename: str) -> dict:
    """HTML: BeautifulSoup text extraction."""
    from bs4 import BeautifulSoup
    import chardet

    detected = chardet.detect(file_bytes)
    encoding = detected.get("encoding") or "utf-8"
    try:
        html = file_bytes.decode(encoding)
    except (UnicodeDecodeError, LookupError):
        html = file_bytes.decode("utf-8", errors="replace")

    soup = BeautifulSoup(html, "lxml")
    # Remove script and style elements
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    text = soup.get_text(separator="\n", strip=True)
    return _result(text, "html_bs4", "html")


def _extract_epub(file_bytes: bytes, filename: str) -> dict:
    """EPUB: ebooklib + BeautifulSoup."""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup

        book = epub.read_epub(io.BytesIO(file_bytes))
        parts = []

        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "lxml")
            text = soup.get_text(separator="\n", strip=True)
            if text.strip():
                parts.append(text)

        return _result("\n\n".join(parts), "epub_ebooklib", "epub")
    except Exception as e:
        logger.warning(f"EPUB extraction failed: {e}")
        return {"text": "", "length": 0, "method": "epub_failed", "format": "epub", "error": str(e)}


def _extract_image_ocr(file_bytes: bytes, filename: str) -> dict:
    """Image OCR via tesseract."""
    try:
        import pytesseract
        from PIL import Image

        img = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(img, lang="ukr+deu+eng+rus+pol")
        return _result(text, "ocr_tesseract", os.path.splitext(filename.lower())[1].lstrip("."))
    except Exception as e:
        logger.warning(f"Image OCR failed: {e}")
        return {"text": "", "length": 0, "method": "ocr_failed",
                "format": os.path.splitext(filename.lower())[1].lstrip("."), "error": str(e)}


def _extract_text_file(file_bytes: bytes, filename: str) -> dict:
    """Plain text with encoding detection."""
    import chardet

    ext = os.path.splitext(filename.lower())[1].lstrip(".")
    detected = chardet.detect(file_bytes)
    encoding = detected.get("encoding") or "utf-8"

    try:
        text = file_bytes.decode(encoding)
    except (UnicodeDecodeError, LookupError):
        text = file_bytes.decode("utf-8", errors="replace")

    if text.startswith("\ufeff"):
        text = text[1:]

    return _result(text, "text_decoded", ext or "txt")


# ─── Helpers ─────────────────────────────────────────────────────────────────

def _result(text: str, method: str, fmt: str, content_type: str = "text") -> dict:
    """Build standardized result with truncation."""
    text = text.strip()
    truncated = False
    max_len = 1_000_000 if content_type == "html" else MAX_TEXT_LENGTH
    if len(text) > max_len:
        if content_type == "html":
            text = text[:max_len] + "\n<p><em>[... текст обрізано ...]</em></p>"
        else:
            text = text[:MAX_TEXT_LENGTH] + "\n\n[... текст обрізано (ліміт 500к символів) ...]"
        truncated = True

    return {
        "text": text,
        "length": len(text),
        "method": method,
        "format": fmt,
        "truncated": truncated,
        "content_type": content_type,
    }
